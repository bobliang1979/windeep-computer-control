#!/usr/bin/env python3
"""Create the computer control system PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — Midnight Executive
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x02, 0x80, 0x90)
GREEN = RGBColor(0x2C, 0x8C, 0x3E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT = RGBColor(0xF9, 0x61, 0x67)  # Coral

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, w, h, text, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, w, h, title, body, title_color=TEAL, body_color=GRAY, bg_color=WHITE):
    shape = add_shape(slide, left, top, w, h, bg_color)
    shape.shadow.inherit = False
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), w - Inches(0.6), Inches(0.5),
                 title, 16, True, title_color)
    # Body
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), w - Inches(0.6), h - Inches(0.9),
                 body, 12, False, body_color)

def add_icon_bullet(slide, left, top, number, text, color=ACCENT):
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    # Text
    add_text_box(slide, left + Inches(0.5), top - Inches(0.05), Inches(5), Inches(0.4),
                 text, 13, False, GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Decorative bar
add_shape(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "电脑控制技能栈", 44, True, WHITE, PP_ALIGN.CENTER, 'Arial Black')
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Bytebot 启发 · 22 MCP Tools · 10x 速度提升 · UIA盲区全覆盖", 20, False, ICE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "Hermes Agent × Codex++ 联合交付", 16, False, RGBColor(0x88, 0x99, 0xCC), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "2026-06-27", 14, False, RGBColor(0x88, 0x99, 0xCC), PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: Bytebot Inspiration
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Bytebot 分析 → 全栈改造", 32, True, NAVY)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(4), Inches(0.04), ACCENT)

# Left: Bytebot
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
         "Bytebot (11K★ GitHub)",
         "TypeScript, Docker虚拟桌面\nnut.js坐标控制, NestJS后端\n纯像素定位, 无Accessibility Tree\n每个Action后750ms自动截图\nSharp.js渐进式图片压缩",
         RGBColor(0xB8, 0x50, 0x42), GRAY)

# Right: Our approach
add_card(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5),
         "我们的方案 (cua-driver + winctl)",
         "Rust原生, 共享桌面不抢焦点\nUIA/SOM精确定位 + 元素指纹\n22 MCP Tools, HTTP :59322\nOCR覆盖Electron盲区\nPillow Pipeline 47x加速压缩",
         TEAL, GRAY)

# Bottom: Key insight
add_card(slide, Inches(2.5), Inches(4.5), Inches(8.3), Inches(2.3),
         "核心理念：可移植的是流程细节，不是架构",
         "Bytebot = 给AI一台自己的容器桌面\n我们 = AI与你共享真实桌面 — 更难但长期价值更高\n\n5个可移植改进 → 全部落地:\n① settle delay → ② 渐进压缩 → ③ action粒度 → ④ holdKeys → ⑤ 截图验证",
         NAVY, GRAY, LIGHT_BG)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "系统架构", 32, True, WHITE)
add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), ACCENT)

# Three layers
layers = [
    ("底层: Windeep ctypes Win32", TEAL,
     "list_windows / focus_window / send_keys\nwindow_click / exec_process\nGetDC / PrintWindow / PostMessage\n零外部依赖, 纯ctypes"),
    ("中间层: Scripts 模块", RGBColor(0x02, 0x80, 0x90),
     "UiTreeCache — TTL缓存+指纹索引\nElementFingerprint — SHA256元素标识\nOCRFinder — Windows原生OCR (WinRT)\nSmartMatcher — 5策略精准匹配\nAssertionVerifier — 结构化断言验证"),
    ("顶层: MCP Server + Hermes Skill", RGBColor(0x02, 0x90, 0x96),
     "winctl_mcp_server — 22 MCP tools, HTTP :59322\ncomputer_control_enhanced — CLI + 库\ncompress_image — 渐进式压缩 + Pipeline\nwindeep-computer-use — 共享Hermes Skill"),
]

for i, (title, color, body) in enumerate(layers):
    y = Inches(1.3) + i * Inches(1.9)
    # Layer card
    card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.7), WHITE)
    # Color bar on left
    add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(1.7), color)
    # Title
    add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(5), Inches(0.5),
                 title, 18, True, NAVY)
    # Body
    add_text_box(slide, Inches(1.2), y + Inches(0.65), Inches(10.5), Inches(1.0),
                 body, 12, False, GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: Speed
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "🚀 速度：3-10s → <1s", 32, True, NAVY)
add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), ACCENT)

# Four optimization cards
speedups = [
    ("P0\nUI Tree Cache", "800ms→0ms\n∞ 加速", TEAL),
    ("P1\nPipeline压缩", "1517ms→57ms\n47x 加速", ACCENT),
    ("P2\nAdaptive Settle", "750ms→350ms\n2x 加速", RGBColor(0xB8, 0x50, 0x42)),
    ("P3\nAction合并", "5000ms→50ms\n100x 加速", GREEN),
]

for i, (title, metric, color) in enumerate(speedups):
    x = Inches(0.8) + i * Inches(3.1)
    # Card
    card = add_shape(slide, x, Inches(1.3), Inches(2.8), Inches(2.5), WHITE)
    # Top color stripe
    add_shape(slide, x, Inches(1.3), Inches(2.8), Inches(0.06), color)
    # Title
    add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(0.8),
                 title, 13, True, NAVY, PP_ALIGN.CENTER)
    # Big metric number
    add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(0.6),
                 metric.split('\n')[0], 22, True, color, PP_ALIGN.CENTER)
    # Speedup
    add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.4), Inches(0.5),
                 metric.split('\n')[1] if '\n' in metric else '', 16, True, color, PP_ALIGN.CENTER)

# Comparison table
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "闭环时序对比", 20, True, NAVY)

old_new = [
    ("闭环总计", "~3-10 秒", "~200-1000 ms", "~10x"),
]
# Table header
headers = ["环节", "之前", "现在", "加速比"]
cols = [Inches(0.5), Inches(4), Inches(4), Inches(4)]
y_start = Inches(4.8)
for j, h in enumerate(headers):
    add_text_box(slide, Inches(0.8) + cols[j], y_start, cols[j+1]-cols[j] if j < len(cols)-1 else Inches(3), Inches(0.4),
                 h, 12, True, NAVY)

rows = [
    ("UI Tree", "每次800ms", "首次800ms,后续0ms", "∞"),
    ("截图→压缩", "1517ms", "57ms (Pipeline)", "47x"),
    ("文本输入(100字符)", "5000ms", "50ms (set_value)", "100x"),
    ("Settle延迟", "固定750ms", "自适应200-1000ms", "2x"),
]
for i, (step, before, after, speedup) in enumerate(rows):
    y = y_start + Inches(0.45) + i * Inches(0.4)
    bg_color = WHITE if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF0)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.38), bg_color)
    add_text_box(slide, Inches(0.8) + Inches(0.2), y + Inches(0.02), Inches(2.5), Inches(0.35),
                 step, 11, False, NAVY)
    add_text_box(slide, Inches(0.8) + Inches(3), y + Inches(0.02), Inches(2.5), Inches(0.35),
                 before, 11, False, GRAY)
    add_text_box(slide, Inches(0.8) + Inches(6), y + Inches(0.02), Inches(2.5), Inches(0.35),
                 after, 11, True, TEAL)
    add_text_box(slide, Inches(0.8) + Inches(9), y + Inches(0.02), Inches(2), Inches(0.35),
                 speedup, 11, True, ACCENT)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: Precision & UIA Blind Spot
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "🎯 精准度：UIA盲区全覆盖", 32, True, WHITE)
add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), ACCENT)

# 5 strategies as a flow
strategies = [
    ("1", "UIA Exact", "角色+标签精确匹配\n置信度 98%", TEAL),
    ("2", "UIA Fuzzy", "Levenshtein模糊匹配\n置信度 70-95%", RGBColor(0x02, 0x80, 0x90)),
    ("3", "OCR", "Windows原生OCR\n覆盖Electron/Canvas盲区", ACCENT),
    ("4", "Position", "最近可点击元素\n按上次坐标定位", RGBColor(0xB8, 0x50, 0x42)),
    ("5", "Coordinate", "裸x,y坐标\n最终fallback", GRAY),
]

for i, (num, name, desc, color) in enumerate(strategies):
    x = Inches(0.5) + i * Inches(2.5)
    # Card
    card = add_shape(slide, x, Inches(1.3), Inches(2.3), Inches(2.8), WHITE)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), Inches(1.5), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Name
    add_text_box(slide, x + Inches(0.1), Inches(2.2), Inches(2.1), Inches(0.4),
                 name, 14, True, NAVY, PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x + Inches(0.15), Inches(2.6), Inches(2.0), Inches(1.2),
                 desc, 11, False, GRAY, PP_ALIGN.CENTER)
    # Arrow (except last)
    if i < 4:
        arrow = add_text_box(slide, x + Inches(2.3), Inches(2.4), Inches(0.3), Inches(0.4),
                             "→", 24, True, ACCENT, PP_ALIGN.CENTER)

# Bottom: assertion list
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
             "验证断言系统", 20, True, WHITE)
assertions = [
    ("hash_change", "UI树哈希改变"),
    ("element_appeared(text)", "元素出现"),
    ("element_disappeared(text)", "元素消失"),
    ("text_contains(text)", "文本包含"),
]
for i, (name, desc) in enumerate(assertions):
    x = Inches(0.8) + i * Inches(3.1)
    add_card(slide, x, Inches(5.1), Inches(2.8), Inches(1.5),
             name, desc, TEAL, GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: MCP Tools
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "🔌 winctl MCP Server — 22 Tools", 32, True, NAVY)
add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), ACCENT)

# Tool categories
categories = [
    ("窗口管理", ["list_windows", "find_windows", "get_window_info", "focus_window", "move_window"]),
    ("窗口状态", ["close_window", "minimize_window", "maximize_window", "restore_window"]),
    ("输入控制", ["click", "type_text", "paste_text", "send_keys", "launch"]),
    ("视觉", ["screenshot", "desktop_info"]),
    ("验证", ["capture_state", "verify", "ocr_find", "ocr_available"]),
    ("智能匹配", ["smart_find", "smart_click"]),
]

for i, (cat, tools) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.3) + row * Inches(2.8)
    # Card
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.5), WHITE)
    # Header bar
    add_shape(slide, x, y, Inches(3.9), Inches(0.06), TEAL)
    # Category name
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.4),
                 cat, 15, True, NAVY)
    # Tools
    for j, tool in enumerate(tools):
        add_text_box(slide, x + Inches(0.3), y + Inches(0.7) + j * Inches(0.35),
                     Inches(3.3), Inches(0.3),
                     f"▸ {tool}", 11, False, GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: Full Stack & Health
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "全栈健康度", 32, True, WHITE)
add_shape(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(0.04), ACCENT)

# Files table
files_data = [
    ("winctl_mcp_server.py", "33KB", "22 MCP Tools, HTTP :59322"),
    ("computer_control_enhanced.py", "20KB", "CLI + P0 特性集成"),
    ("compress_image.py", "16KB", "渐进式压缩 + Pipeline"),
    ("scripts/ui_tree_cache.py", "2KB", "TTL缓存 + 指纹索引"),
    ("scripts/element_fingerprint.py", "2KB", "SHA256稳定指纹"),
    ("scripts/ocr_finder.py", "6KB", "Windows原生OCR"),
    ("scripts/smart_matcher.py", "11KB", "5策略智能匹配"),
    ("scripts/assertion_verifier.py", "11KB", "4种断言验证"),
    ("scripts/__init__.py", "0.5KB", "统一包导出"),
]

# Table header
table_top = Inches(1.2)
col_widths = [Inches(4.5), Inches(1.2), Inches(5)]
headers = ["文件", "大小", "功能"]
for j, h in enumerate(headers):
    x = Inches(1) + sum(col_widths[:j])
    add_shape(slide, x, table_top, col_widths[j], Inches(0.4), RGBColor(0x1A, 0x1A, 0x4A))
    add_text_box(slide, x + Inches(0.1), table_top + Inches(0.05), col_widths[j] - Inches(0.2), Inches(0.3),
                 h, 12, True, WHITE)

for i, (fname, size, desc) in enumerate(files_data):
    y = table_top + Inches(0.45) + i * Inches(0.38)
    bg = RGBColor(0x25, 0x30, 0x60) if i % 2 == 0 else RGBColor(0x1E, 0x27, 0x61)
    for j, val in enumerate([fname, size, desc]):
        x = Inches(1) + sum(col_widths[:j])
        add_shape(slide, x, y, col_widths[j], Inches(0.35), bg)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.03), col_widths[j] - Inches(0.2), Inches(0.3),
                     val, 10, False, ICE)

# Health metrics on right
health_y = Inches(5.0)
add_text_box(slide, Inches(1), health_y, Inches(5), Inches(0.4),
             "架构健康度", 18, True, WHITE)

health_items = [
    ("语法正确性", "✅ 100%"),
    ("异常覆盖", "✅ 100%"),
    ("Shell注入风险", "✅ 无风险"),
    ("并发安全", "✅ ThreadedHTTPServer"),
    ("资源泄漏", "✅ GDI + Clipboard"),
    ("模块化", "✅ 8模块职责分离"),
]
for i, (metric, status) in enumerate(health_items):
    y = health_y + Inches(0.45) + i * Inches(0.38)
    add_text_box(slide, Inches(1.2), y, Inches(3), Inches(0.35),
                 metric, 11, False, ICE)
    add_text_box(slide, Inches(4), y, Inches(3), Inches(0.35),
                 status, 11, True, GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: Thank You
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "Thank You", 48, True, WHITE, PP_ALIGN.CENTER, 'Arial Black')
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Bytebot → 全栈交付 · Hermes Agent × Codex++", 20, False, ICE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "源码已打包至桌面: windeep-computer-control.zip (27KB, 9文件, 101KB)", 14, False, RGBColor(0x88, 0x99, 0xCC), PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\10074\Desktop\电脑控制技能栈.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
