#!/usr/bin/env python3
"""Generate 王小慧艺术馆 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = RGBColor(255, 255, 255)
G = RGBColor(200, 200, 200)
BG = RGBColor(30, 30, 30)
ACCENT = RGBColor(210, 180, 140)  # Gold
RED = RGBColor(180, 60, 60)

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=W, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet_slide(slide, items, left=1.0, top=2.0, font_size=20):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = G
        p.space_after = Pt(8)

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.5, "王小慧艺术馆", 54, W, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.5, 11, 1, "Wang Xiaohui Art Museum", 28, ACCENT, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.8, "艺术 · 科技 · 生活", 22, G, False, PP_ALIGN.CENTER)

# ── Slide 2: 艺术家简介 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 1, "艺术家简介 / About the Artist", 36, ACCENT, True)
items = [
    "▸ 王小慧，旅德艺术家，跨界创作先锋",
    "▸ 毕业于上海同济大学建筑系，后在慕尼黑美术学院深造",
    "▸ 作品涵盖摄影、影像、装置、设计、写作等多个领域",
    "▸ 曾获世界杰出华人艺术家、德国国家艺术奖等荣誉",
    "▸ 作品被世界多家艺术馆收藏，在国际上具有广泛影响力",
    "",
    "▸ Wang Xiaohui, Chinese-German artist, pioneer of跨界 creation",
    "▸ Graduated from Tongji University, studied at Academy of Fine Arts Munich",
    "▸ Works span photography, video, installation, design, and writing",
]
add_bullet_slide(slide, items, font_size=18)

# ── Slide 3: 艺术馆简介 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 1, "艺术馆简介 / Museum Overview", 36, ACCENT, True)
items = [
    "▸ 王小慧艺术馆位于中国上海，是展现其艺术理念的专属空间",
    "▸ 建筑面积约 3000 平方米，融合展览、创作、交流等多功能",
    "▸ 由著名建筑师设计，建筑本身即是一件艺术品",
    "▸ 常设展览：王小慧各时期代表作品",
    "▸ 特展空间：每年举办 4-6 场国内外当代艺术展览",
    "▸ 公共教育：工作坊、讲座、导览等公益活动",
]
add_bullet_slide(slide, items, font_size=18)

# ── Slide 4: 艺术特点 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 1, "艺术特点 / Artistic Features", 36, ACCENT, True)
items = [
    "▸ 跨界融合：打破艺术门类界限，创造新的视觉语言",
    "▸ 科技结合：运用数字媒体、影像技术拓展艺术表达",
    "▸ 文化对话：中西文化交汇，传统与现代的碰撞",
    "▸ 生命主题：关注生命、爱情、自然等永恒命题",
    "▸ 互动体验：观众参与成为作品的一部分",
]
add_bullet_slide(slide, items, font_size=18)

# ── Slide 5: 代表作品 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 1, "代表作品 / Representative Works", 36, ACCENT, True)
works = [
    ("《莲之梦》系列", "摄影装置，通过荷花意象表现生命轮回"),
    ("《女性之声》", "多媒体装置，探讨女性在社会中的角色"),
    ("《时间的碎片》", "影像作品，以碎片化叙事展现记忆与时间"),
    ("《花之灵》", "数字艺术，花卉与光影的数字化演绎"),
    ("《桥》系列", "公共艺术装置，连接不同文化与空间"),
]
y = 1.8
for title, desc in works:
    add_text_box(slide, 1.0, y, 11, 0.5, title, 22, ACCENT, True)
    add_text_box(slide, 1.5, y + 0.5, 10, 0.4, desc, 16, G)
    y += 1.2

# ── Slide 6: 参观信息 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 1, "参观信息 / Visit Information", 36, ACCENT, True)
items = [
    "📍 地址：上海市（具体地址详见官网）",
    "🕐 开放时间：周二至周日 10:00-18:00（周一闭馆）",
    "🎫 门票：成人票 80 元，学生/老人 40 元",
    "🔗 官网：www.wangxiaohui-artmuseum.com",
    "📞 电话：+86-21-xxxx-xxxx",
    "",
    "📍 Address: Shanghai, China (see website for details)",
    "🕐 Hours: Tue-Sun 10:00-18:00 (closed Mon)",
    "🎫 Ticket: Adult 80 RMB, Student/Senior 40 RMB",
]
add_bullet_slide(slide, items, font_size=18)

# ── Slide 7: End ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 2.0, 11, 1, "感谢观看", 48, W, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.5, 11, 1, "Thank You", 32, ACCENT, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.0, 11, 0.6, "欢迎莅临王小慧艺术馆", 20, G, False, PP_ALIGN.CENTER)

# Save
path = r"C:\Users\10074\Desktop\王小慧艺术馆.pptx"
prs.save(path)
print(f"PPT saved: {path}")
print(f"Size: {os.path.getsize(path)/1024:.0f} KB")
print(f"Slides: {len(prs.slides)}")
